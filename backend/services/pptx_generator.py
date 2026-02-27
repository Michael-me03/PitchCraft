"""
PitchCraft PPTX Generator — McKinsey-Style Layout Engine
========================================================

Design principles
-----------------
• ALL positions are derived from the actual slide dimensions (SlideGeometry).
  The layout works identically on 4:3, 16:9 and any custom template size.
• Template placeholders are NEVER used for positioning.  They are removed and
  every element is drawn manually in fixed grid zones → zero overlap guaranteed.
• Three exclusive vertical zones per slide:
    ┌─────────────────────────────────────────────────────────┐
    │▓▓▓▓▓▓ top accent bar  0 % – 1.2 %  ▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓│
    │  TITLE ZONE           5 % – 18 %                       │
    │──────────────────────────────────── divider  19 % ─────│
    │  CONTENT ZONE        21 % – 89.5 %                     │
    │                                  slide-nr  93 % – 97 % │
    └─────────────────────────────────────────────────────────┘
• Font sizes are adaptive: computed from text length to prevent overflow.
• McKinsey color palette: deep navy primary, controlled accent colors.
"""

import io
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from models.schemas import PresentationStructure, SlideContent, ChartSpec
from services.chart_engine import render_chart, set_chart_theme

logger = logging.getLogger(__name__)


# ─── Brand Font ───────────────────────────────────────────────────────────────

FONT = "Equip Medium"   # Nagarro brand font — embedded in Nagarro template


# ─── Default Color Palette (overridden per-render by template colors) ─────────

_DEFAULT_COLORS: dict = {
    "bg":         RGBColor(0xF8, 0xF9, 0xFA),
    "navy":       RGBColor(0x00, 0x27, 0x5A),
    "blue":       RGBColor(0x00, 0x5B, 0xB5),
    "teal":       RGBColor(0x00, 0x96, 0xA8),
    "green":      RGBColor(0x00, 0x7A, 0x4C),
    "red":        RGBColor(0xCC, 0x00, 0x00),
    "orange":     RGBColor(0xE8, 0x7B, 0x1E),
    "text_dark":  RGBColor(0x1A, 0x1A, 0x2E),
    "text_body":  RGBColor(0x2D, 0x3A, 0x4E),
    "text_muted": RGBColor(0x64, 0x74, 0x8B),
    "divider":    RGBColor(0xCC, 0xD1, 0xD9),
    "bg_card":    RGBColor(0xF5, 0xF7, 0xFA),
    "bg_accent":  RGBColor(0xEE, 0xF4, 0xFF),
}

# Active palette — replaced at the start of every generate_pptx() call
COLORS: dict = _DEFAULT_COLORS


# ─── Template Introspection Helpers ──────────────────────────────────────────

def _is_dark_template() -> bool:
    """
    Return True when the active template has a dark slide background.

    Checks COLORS['bg'] luminance directly.  RGBColor is a (r, g, b) tuple.
    On dark templates luminance < 128; on light templates it is near-white.
    """
    c = COLORS["bg"]
    r, g, b = c[0], c[1], c[2]
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) < 128


# ─── Template Color Helpers ────────────────────────────────────────────────────

def _hex_to_rgb_color(hex_str: str) -> RGBColor:
    """Convert a CSS hex color string to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _is_dark_color(hex_str: str) -> bool:
    """Return True if the color is perceptually dark (luminance < 50%)."""
    h  = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) < 128


def _build_template_palette(tc: dict) -> dict:
    """
    Derive a full slide color palette from template accent/text/muted/bg values.

    On dark backgrounds: card fills become dark-tinted, text uses the template's
    light text color.  On light backgrounds: existing neutral card fills are kept.

    Args:
        tc: Template colors dict with keys bg, accent, text, muted (hex strings).

    Returns:
        New COLORS-compatible dict with RGBColor values.
    """
    bg_hex     = tc.get("bg",     "#FFFFFF")
    accent_hex = tc.get("accent", "#0027A0")
    text_hex   = tc.get("text",   "#1A1A2E")
    muted_hex  = tc.get("muted",  "#64748B")

    dark_bg    = _is_dark_color(bg_hex)
    accent_rgb = _hex_to_rgb_color(accent_hex)
    muted_rgb  = _hex_to_rgb_color(muted_hex)

    # Auto-derive text colour from background so slides are always readable,
    # regardless of what the template config says.
    if dark_bg:
        text_rgb = RGBColor(0xFF, 0xFF, 0xFF)          # white on dark bg
    else:
        text_rgb = RGBColor(0x0F, 0x17, 0x2A)          # near-black on light bg

    if dark_bg:
        bg_r = int(bg_hex.lstrip("#")[0:2], 16)
        bg_g = int(bg_hex.lstrip("#")[2:4], 16)
        bg_b = int(bg_hex.lstrip("#")[4:6], 16)
        card_rgb    = RGBColor(min(255, bg_r + 28), min(255, bg_g + 28), min(255, bg_b + 28))
        divider_rgb = RGBColor(min(255, bg_r + 55), min(255, bg_g + 55), min(255, bg_b + 55))
    else:
        card_rgb    = _DEFAULT_COLORS["bg_card"]
        divider_rgb = _DEFAULT_COLORS["divider"]

    return {
        "bg":         _hex_to_rgb_color(bg_hex),
        "navy":       accent_rgb,
        "blue":       accent_rgb,
        "teal":       accent_rgb,
        "green":      _DEFAULT_COLORS["green"],
        "red":        _DEFAULT_COLORS["red"],
        "orange":     _DEFAULT_COLORS["orange"],
        "text_dark":  text_rgb,
        "text_body":  text_rgb,
        "text_muted": muted_rgb,
        "divider":    divider_rgb,
        "bg_card":    card_rgb,
        "bg_accent":  card_rgb,
    }


# ─── Layout Grid  (fractions of slide width W / height H) ────────────────────

ML       = 0.050   # left margin
CW       = 0.900   # content width  (ML … 0.950)

BAR_H    = 0.012   # top accent bar height
TITLE_Y  = 0.050   # title zone top
TITLE_H  = 0.130   # title zone height  → bottom 0.180
RULE_Y   = 0.190   # thin divider rule
BODY_Y   = 0.210   # content zone top
BODY_H   = 0.685   # content zone height → bottom 0.895
FOOTER_Y = 0.930   # slide-number row


# ─── Logo-Safe Zone Detection ─────────────────────────────────────────────────

def _detect_logo_safe_y(prs: Presentation) -> float:
    """
    Scan slide layouts for image shapes in the top-left logo zone.

    Returns the minimum Y fraction (0–0.30) below which it is safe to place the
    slide title, so that PitchCraft text never overlaps template logos.

    Args:
        prs: Loaded Presentation object.

    Returns:
        Safe top-Y fraction (0.0 for templates without top-left logos).
    """
    W      = prs.slide_width
    H      = prs.slide_height
    safe_y = 0.0

    shape_lists = [prs.slide_master.shapes] + [
        lay.shapes for lay in prs.slide_master.slide_layouts
    ]
    for shapes in shape_lists:
        for sh in shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            t = (sh.top    or 0) / H
            b = t + (sh.height or 0) / H
            l = (sh.left   or 0) / W
            # Logo heuristic: starts in top 15 % AND within left 30 % of slide
            if t < 0.15 and l < 0.30:
                safe_y = max(safe_y, b + 0.03)   # 3 % breathing room

    return min(safe_y, 0.30)   # never push title below 30 %


# ─── Slide Geometry ───────────────────────────────────────────────────────────

class SlideGeometry:
    """
    Converts layout-grid fractions to absolute EMU positions.

    Zone fractions are computed dynamically: if a template logo occupies the
    top-left corner, the title zone shifts down to avoid overlap.
    """

    def __init__(self, prs: Presentation, logo_safe_y: float = 0.0):
        self.W = prs.slide_width
        self.H = prs.slide_height

        # ── Dynamic vertical zones ─────────────────────────────────────────────
        self.title_y = max(TITLE_Y, logo_safe_y)          # top of title zone
        self.title_h = TITLE_H                             # height of title zone
        self.rule_y  = self.title_y + self.title_h + 0.01 # thin divider
        self.body_y  = self.rule_y  + 0.020               # content zone top
        self.body_h  = FOOTER_Y - 0.025 - self.body_y     # content zone height

    def x(self, pct: float) -> int:
        return int(self.W * pct)

    def y(self, pct: float) -> int:
        return int(self.H * pct)

    def w(self, pct: float) -> int:
        return int(self.W * pct)

    def h(self, pct: float) -> int:
        return int(self.H * pct)


# ─── Low-level Drawing Primitives ─────────────────────────────────────────────

def _remove_all_placeholders(slide):
    """Remove every template placeholder so we position everything ourselves."""
    to_remove = [ph._element for ph in slide.placeholders]
    for el in to_remove:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def _rect(slide, left, top, width, height, fill: RGBColor):
    """Plain rectangle with solid fill, no visible border."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _round_rect(slide, left, top, width, height, fill: RGBColor):
    """Rounded rectangle with solid fill, no border."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height, word_wrap: bool = True):
    """Add a textbox with word-wrap enabled."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = word_wrap
    return txb


def _set_para(p, text: str, size: int, bold: bool = False,
              color: RGBColor = None, align=PP_ALIGN.LEFT):
    """Set text, size, bold, color, font and alignment on a paragraph."""
    p.text = text
    p.font.name  = FONT
    p.font.size  = Pt(size)
    p.font.bold  = bold
    p.font.color.rgb = color or COLORS["text_dark"]
    p.alignment  = align


def _vcenter(txb) -> None:
    """Vertically center text within a textbox via bodyPr anchor attribute."""
    txb.text_frame._txBody.bodyPr.set("anchor", "ctr")


def _set_line_spacing(p, spacing: float = 1.35):
    """Inject line-spacing XML into a paragraph."""
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
    pPr.append(lnSpc)
    spcPct = lnSpc.makeelement(qn("a:spcPct"), {})
    spcPct.set("val", str(int(spacing * 100_000)))
    lnSpc.append(spcPct)


def _enable_auto_shrink(txb, min_scale: int = 40) -> None:
    """Enable PowerPoint's native text auto-shrink so text never overflows the shape.

    Injects ``<a:normAutofit fontScale="..." lnSpcReduction="20000"/>`` into the
    text frame's ``<a:bodyPr>`` element.  PowerPoint will progressively reduce the
    font size (down to *min_scale* percent of the original) and tighten line
    spacing until the text fits within the shape boundaries.

    Args:
        txb:       Shape (textbox or auto-shape) whose text frame to configure.
        min_scale: Minimum font-scale percentage (100000 = 100 %).  Lower values
                   allow more aggressive shrinking.  50 (= 50 %) is a safe floor.
    """
    bodyPr = txb.text_frame._txBody.bodyPr
    # Remove any existing autofit / spAutoFit so normAutofit takes effect
    for child in list(bodyPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("normAutofit", "spAutoFit", "noAutofit"):
            bodyPr.remove(child)
    autofit = bodyPr.makeelement(qn("a:normAutofit"), {})
    autofit.set("fontScale", str(min_scale * 1000))   # e.g. 50 → "50000"
    autofit.set("lnSpcReduction", "20000")             # allow 20 % line-spacing reduction
    bodyPr.append(autofit)


# ─── Adaptive Typography ──────────────────────────────────────────────────────

def _title_size(text: str, base: int = 32) -> int:
    """Scale font size down for long titles to prevent overflow."""
    n = len(text)
    if n > 90: return max(base - 10, 20)
    if n > 70: return max(base - 6,  24)
    if n > 50: return max(base - 4,  26)
    if n > 35: return max(base - 2,  28)
    return base


def _bullet_size(bullets: list) -> int:
    """Adaptive bullet font size: fewer / shorter bullets → larger text.

    Conservative sizing to prevent text from overflowing constrained shapes.
    Uses aggressive down-scaling for long or numerous bullets.
    """
    if not bullets:
        return 14
    n    = len(bullets)
    avg  = sum(len(b) for b in bullets) / n
    maxl = max(len(b) for b in bullets)
    total = sum(len(b) for b in bullets)
    # Very long individual bullets or high total volume → force small text
    if maxl > 120 or total > 500 or (n >= 4 and avg > 70):
        return 11
    if maxl > 100 or total > 400 or (n >= 4 and avg > 60):
        return 12
    if maxl > 80 or (n >= 3 and avg > 50):
        return 13
    if n <= 3 and avg <= 35:  return 18
    if n <= 3 and avg <= 50:  return 16
    if n <= 4 and avg <= 45:  return 15
    if n <= 5:                return 14
    if n <= 7:                return 13
    return 12


# ─── Frame Elements (shared by all content slides) ───────────────────────────

def _add_top_bar(slide, g: SlideGeometry):
    """Full-width navy accent bar at the very top of the slide.

    On dark templates the blank slide layout often inherits a light/white
    background from the source PPTX.  Paint the full-slide background first
    so dark-text templates don't end up with invisible white-on-white text.
    """
    if _is_dark_template():
        _rect(slide, 0, 0, g.W, g.H, COLORS["bg"])
    _rect(slide, 0, 0, g.W, g.h(BAR_H), COLORS["navy"])


def _add_divider(slide, g: SlideGeometry):
    """Thin horizontal rule that separates title zone from content zone."""
    _rect(
        slide,
        g.x(ML), g.y(g.rule_y),
        g.w(CW),  g.h(0.002),
        COLORS["divider"],
    )


def _add_title(slide, g: SlideGeometry, text: str, base_size: int = 32):
    """Draw the slide title at the dynamic title_y grid position."""
    size = _title_size(text, base_size)
    txb  = _textbox(slide, g.x(ML), g.y(g.title_y), g.w(CW), g.h(g.title_h))
    _vcenter(txb)
    _enable_auto_shrink(txb)
    p    = txb.text_frame.paragraphs[0]
    _set_para(p, text, size, bold=True, color=COLORS["text_dark"])
    return txb


def _add_slide_number(slide, g: SlideGeometry, num: int, total: int):
    """Slide number in the bottom-right corner."""
    txb = _textbox(slide, g.x(0.80), g.y(FOOTER_Y), g.w(0.15), g.h(0.05))
    p   = txb.text_frame.paragraphs[0]
    _set_para(p, f"{num} / {total}", 8,
              color=COLORS["text_muted"], align=PP_ALIGN.RIGHT)


def _slide_frame(slide, g: SlideGeometry, title: str, base_size: int = 28):
    """Draw the standard frame: top bar + title + divider rule."""
    _add_top_bar(slide, g)
    _add_title(slide, g, title, base_size)
    _add_divider(slide, g)


# ─── Bullet List Helper ───────────────────────────────────────────────────────

def _add_bullets(tf, bullets: list, size: int = 14,
                 bold_first_word: bool = True,
                 color: RGBColor = None,
                 line_spacing: float = 1.15):
    """
    Render a bullet list into a text frame.
    McKinsey style: em-dash prefix, bold keyword before colon.
    """
    tf.clear()
    tf.word_wrap = True
    body_color = color or COLORS["text_body"]
    EM = "\u2014 "  # em-dash prefix

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after  = Pt(6)
        p.level        = 0
        p.alignment    = PP_ALIGN.LEFT
        _set_line_spacing(p, line_spacing)

        if bold_first_word and ":" in bullet:
            label, rest = bullet.split(":", 1)
            r1 = p.add_run()
            r1.text           = EM + label + ":"
            r1.font.name      = FONT
            r1.font.size      = Pt(size)
            r1.font.bold      = True
            r1.font.color.rgb = COLORS["text_dark"]
            r2 = p.add_run()
            r2.text           = rest
            r2.font.name      = FONT
            r2.font.size      = Pt(size)
            r2.font.color.rgb = body_color
        else:
            r = p.add_run()
            r.text           = EM + bullet
            r.font.name      = FONT
            r.font.size      = Pt(size)
            r.font.color.rgb = body_color


# ─── Content Layout Helpers ───────────────────────────────────────────────────

def _has_emoji_prefix(bullets: list) -> bool:
    """True when every bullet starts with an emoji / high-codepoint character."""
    return all(b and ord(b[0]) > 0x1F00 for b in bullets)


def _use_card_layout(bullets: list) -> bool:
    """
    Cards only for exactly 2–3 VERY short bullets (pillars / options).
    Long text or emoji-prefixed lists go to icon-list or plain-list instead.
    """
    n = len(bullets)
    if n < 2 or n > 3:
        return False
    if _has_emoji_prefix(bullets):
        return False          # emoji bullets → icon list looks better
    avg_len = sum(len(b) for b in bullets) / n
    return avg_len <= 60


def _build_bullet_cards(slide, g: SlideGeometry, bullets: list) -> None:
    """
    Render 2–3 short bullets as side-by-side rounded cards.
    Each card: coloured accent stripe at top, full bullet text centred below.
    Works on both light and dark templates.

    Args:
        slide:   Target slide object.
        g:       SlideGeometry for this presentation.
        bullets: 2–3 short bullet strings.
    """
    n     = len(bullets)
    bl    = g.x(ML)
    bt    = g.y(g.body_y)
    bw    = g.w(CW)
    bh    = g.h(g.body_h)
    gap_x = g.w(0.025)
    cw    = (bw - gap_x * (n - 1)) // n
    stripe_h = g.h(0.025)

    for i, bullet in enumerate(bullets):
        cx = bl + i * (cw + gap_x)

        # Card background — must be visible on both light and dark slides
        _round_rect(slide, cx, bt, cw, bh, COLORS["bg_card"])

        # Full-width accent top stripe
        _rect(slide, cx, bt, cw, stripe_h, COLORS["navy"])

        # Full bullet text (including any emoji) centred in card
        inner_x = cx + g.w(0.018)
        inner_w = cw - g.w(0.036)
        tx      = _textbox(slide, inner_x, bt + stripe_h + g.h(0.015),
                           inner_w, bh - stripe_h - g.h(0.030))
        _vcenter(tx)
        _enable_auto_shrink(tx)
        tx.text_frame.word_wrap = True
        p = tx.text_frame.paragraphs[0]
        p.font.name      = FONT
        p.font.size      = Pt(18 if n == 2 else 16)
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment      = PP_ALIGN.CENTER
        p.text           = bullet


def _build_icon_list(slide, g: SlideGeometry, bullets: list) -> None:
    """
    Render emoji-prefixed bullets as icon rows: large emoji pill on the left,
    text on the right.  Gives visual variety without full card boxes.

    Each row: coloured rounded pill (emoji) | bold keyword : supporting text

    Args:
        slide:   Target slide object.
        g:       SlideGeometry for this presentation.
        bullets: 2–5 bullet strings that start with an emoji.
    """
    n      = len(bullets)
    bl     = g.x(ML)
    bt     = g.y(g.body_y)
    bw     = g.w(CW)
    bh     = g.h(g.body_h)
    row_h  = bh // n
    pad_y  = g.h(0.010)
    pill_w = g.w(0.065)
    gap    = g.w(0.020)

    for i, bullet in enumerate(bullets):
        ry = bt + i * row_h + pad_y
        rh = row_h - 2 * pad_y

        # Split emoji from body text
        if bullet and ord(bullet[0]) > 0x1F00:
            parts  = bullet.split(" ", 1)
            emoji  = parts[0]
            body   = parts[1].strip() if len(parts) > 1 else ""
        else:
            emoji = "•"
            body  = bullet

        # Coloured circle pill with emoji
        _round_rect(slide, bl, ry + (rh - pill_w) // 2,
                    pill_w, pill_w, COLORS["navy"])
        em_txb = _textbox(slide, bl, ry, pill_w, rh)
        _vcenter(em_txb)
        em_p   = em_txb.text_frame.paragraphs[0]
        em_p.text      = emoji
        em_p.font.size = Pt(22)
        em_p.alignment = PP_ALIGN.CENTER

        # Body text
        tx_x = bl + pill_w + gap
        tx_w = bw - pill_w - gap
        size = 18 if n <= 3 else 16
        body_txb = _textbox(slide, tx_x, ry, tx_w, rh)
        _vcenter(body_txb)
        _enable_auto_shrink(body_txb)
        body_txb.text_frame.word_wrap = True
        p = body_txb.text_frame.paragraphs[0]
        p.font.name = FONT
        p.alignment = PP_ALIGN.LEFT

        if ":" in body:
            kw, rest = body.split(":", 1)
            r1 = p.add_run(); r1.text = kw + ":"; r1.font.name = FONT
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = COLORS["navy"]
            r2 = p.add_run(); r2.text = rest; r2.font.name = FONT
            r2.font.size = Pt(size); r2.font.color.rgb = COLORS["text_body"]
        else:
            r = p.add_run(); r.text = body; r.font.name = FONT
            r.font.size = Pt(size); r.font.color.rgb = COLORS["text_body"]


# ─── Chart Image Helpers ──────────────────────────────────────────────────────

_EMU_PER_INCH = 914_400
_RENDER_DPI   = 300   # 2× HiDPI base for crisp text in all slot sizes


def _compute_render_dims(max_w: int, max_h: int) -> tuple[int, int]:
    """
    Convert slot EMU dimensions to slot-proportional pixel dimensions.

    Caps proportionally so the rendered image always matches the slot aspect
    ratio — capping rw/rh independently would distort the chart when
    python-pptx stretches the image to fill the fixed max_w × max_h slot.

    Returns:
        (render_width, render_height) in pixels, clamped to [400×280, 1920×1080].
    """
    inches_w   = max_w / _EMU_PER_INCH
    inches_h   = max_h / _EMU_PER_INCH
    slot_ratio = inches_w / inches_h if inches_h > 0 else 16 / 9

    rw_cand = min(1920, round(inches_w * _RENDER_DPI))
    rh_cand = round(rw_cand / slot_ratio)
    if rh_cand > 1080:          # too tall → scale down from height
        rh_cand = 1080
        rw_cand = round(rh_cand * slot_ratio)

    return max(400, rw_cand), max(280, rh_cand)


def _render_chart(slide, chart_spec: ChartSpec,
                  left: int, top: int, max_w: int, max_h: int) -> None:
    """
    Render a single chart PNG at slot-proportional dimensions and place it on the slide.

    Uses 300 DPI scaled to the actual slot size so labels and tick marks are
    legible at every grid density — from a full-body chart to a 3×2 dashboard.
    Falls back to an error text box if rendering fails (never crashes the deck).
    """
    rw, rh = _compute_render_dims(max_w, max_h)
    try:
        img_bytes = render_chart(chart_spec.chart_function, chart_spec.params, rw, rh)
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, max_w, max_h)
    except Exception as exc:
        logger.error("Chart render failed [%s]: %s", chart_spec.chart_function, exc)
        txb = _textbox(slide, left, top, max_w, max_h)
        txb.text_frame.text = f"[Chart error: {chart_spec.chart_function}]"


def _render_charts_parallel(
    chart_specs: list[ChartSpec],
    slots: list[tuple[int, int, int, int]],
) -> list[bytes | None]:
    """
    Render multiple charts concurrently using a thread pool.

    chart_engine.render_chart() is thread-safe: it writes render dimensions to
    threading.local() (_ctx), so each worker thread has its own isolated context.

    Args:
        chart_specs: Chart specifications to render.
        slots:       Corresponding (left, top, width, height) tuples in EMUs.

    Returns:
        Index-aligned list of PNG bytes, or None for any chart that failed.
    """
    n       = len(chart_specs)
    results: list[bytes | None] = [None] * n

    def _worker(idx: int, spec: ChartSpec,
                slot: tuple[int, int, int, int]) -> tuple[int, bytes | None]:
        _, _, max_w, max_h = slot
        rw, rh = _compute_render_dims(max_w, max_h)
        try:
            return idx, render_chart(spec.chart_function, spec.params, rw, rh)
        except Exception as exc:
            logger.error("Chart render failed [%s]: %s", spec.chart_function, exc)
            return idx, None

    if n == 1:
        results[0] = _worker(0, chart_specs[0], slots[0])[1]
    else:
        with ThreadPoolExecutor(max_workers=min(4, n)) as executor:
            future_map = {
                executor.submit(_worker, i, spec, slot): i
                for i, (spec, slot) in enumerate(zip(chart_specs, slots))
            }
            for future in as_completed(future_map):
                idx, img = future.result()
                results[idx] = img

    return results


# ─── Layout Discovery ─────────────────────────────────────────────────────────

def _discover_layouts(prs: Presentation) -> dict:
    """
    Discover and map all slide layouts by name patterns.

    Detects both generic PowerPoint layout names (Blank, Title Slide, …) and
    the rich Nagarro / corporate named layouts so that PitchCraft can use each
    template's native designs instead of always falling back to the blank slide.
    """
    layouts: dict = {}

    def _sw(name: str, prefix: str) -> bool:
        """True when name starts with prefix (case-insensitive substring anchor)."""
        return name.startswith(prefix)

    for layout in prs.slide_layouts:
        name     = layout.name.lower()
        ph_count = len(list(layout.placeholders))

        # Numbered Nagarro-style layouts use ONLY prefix matching to avoid
        # cross-matching (e.g. "title blue" is a substring of "2_2_section big title blue").
        # Generic/non-numbered layouts use descriptive substring matching as fallback.

        # ── Title slides ──────────────────────────────────────────────────────
        if _sw(name, "1_1_"):
            layouts.setdefault("title_white", layout)
            layouts.setdefault("title", layout)
        elif _sw(name, "1_2_") or _sw(name, "1_3_"):
            layouts.setdefault("title_blue", layout)
            layouts.setdefault("title", layout)
        elif any(k in name for k in ["title slide", "titelfolie", "titel-"]):
            layouts.setdefault("title", layout)

        # ── Section / divider slides ───────────────────────────────────────────
        elif _sw(name, "2_2_"):
            layouts.setdefault("section_blue", layout)
        elif _sw(name, "2_3_"):
            layouts.setdefault("section_green", layout)
        elif _sw(name, "2_1_"):
            layouts.setdefault("section_white", layout)
        elif _sw(name, "3_2_"):
            layouts.setdefault("section_small_blue", layout)
        elif _sw(name, "3_3_"):
            layouts.setdefault("section_small_green", layout)
        # Generic section divider for non-Nagarro templates
        elif "section" in name and ph_count <= 3 and "photo" not in name:
            layouts.setdefault("section_white", layout)

        # ── Content / bullets ──────────────────────────────────────────────────
        elif _sw(name, "4_2_"):
            layouts.setdefault("info_bullets", layout)
        elif _sw(name, "4_3_"):
            layouts.setdefault("info_numbers", layout)
        elif _sw(name, "5_1_") and ph_count >= 3:
            layouts.setdefault("plain_text", layout)

        # ── Two-column / compare ───────────────────────────────────────────────
        elif _sw(name, "5_3_") or "text in 2 columns" in name:
            layouts.setdefault("two_column_native", layout)
        elif _sw(name, "11_1_compare") or (name.startswith("11_") and "compare" in name):
            layouts.setdefault("compare_native", layout)

        # ── Statement / quote ──────────────────────────────────────────────────
        elif _sw(name, "7_1_") or name == "statement":
            layouts.setdefault("statement", layout)
        elif _sw(name, "12_1_"):
            layouts.setdefault("quote_blue", layout)
        elif _sw(name, "12_2_"):
            layouts.setdefault("quote_purple", layout)

        # ── Closing / thank-you ────────────────────────────────────────────────
        elif _sw(name, "14_2_") or name == "thank you blue":
            layouts.setdefault("closing_blue", layout)
        elif _sw(name, "14_1_") or name == "thank you white":
            layouts.setdefault("closing_white", layout)

        # ── Generic fallbacks ──────────────────────────────────────────────────
        elif any(k in name for k in ["blank", "leer"]):
            layouts.setdefault("blank", layout)
        elif any(k in name for k in ["title only", "nur titel"]):
            layouts.setdefault("title_only", layout)
        elif ph_count >= 2:
            layouts.setdefault("content", layout)

    all_layouts = list(prs.slide_layouts)
    layouts.setdefault("title",      all_layouts[0])
    layouts.setdefault("content",    all_layouts[min(1, len(all_layouts) - 1)])
    layouts.setdefault("blank",      layouts.get("content"))
    layouts.setdefault("title_only", layouts.get("blank"))

    return layouts


_FOOTER_PLACEHOLDER_IDX = 10   # Nagarro (and many corp. templates) use idx 10 as
                               # the repeating footer bar — skip it when filling body


def _fill_native_placeholders(
    slide,
    title: str,
    body: Optional[list] = None,
    body2: Optional[list] = None,
) -> bool:
    """
    Fill a native layout slide's placeholders with content.

    Deliberately does NOT override font name, size, or color so that the
    template's own theme styling (e.g. Nagarro's blue/green sections with
    white text) is preserved exactly.

    Handles non-standard placeholder indices (Nagarro uses idx 10-38 for body
    placeholders) by sorting content placeholders top-to-bottom / left-to-right
    and filling the first two available body slots.  The repeating footer
    placeholder (idx = _FOOTER_PLACEHOLDER_IDX) is always skipped.

    Args:
        slide:  Slide added with a named native layout.
        title:  Text for the TITLE-type placeholder (idx 0).
        body:   Lines for the first content placeholder (top-left body).
        body2:  Lines for the second content placeholder (right / lower body).

    Returns:
        True if at least the title placeholder was found and filled.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    all_phs = list(slide.placeholders)
    if not all_phs:
        return False

    # Separate title placeholder from body/content ones.
    # NOTE: idx=10 is the Nagarro footer bar on most layouts — we exclude it
    # here, but if it turns out to be the *only* text slot (e.g. Statement), we
    # add it back below so that attribution/body text has somewhere to go.
    title_ph        = None
    body_phs        = []
    footer_ph_slots = []   # idx=10 candidates, added back if no other body phs

    for ph in all_phs:
        idx  = ph.placeholder_format.idx
        kind = ph.placeholder_format.type
        if idx == 0 or kind in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title_ph = ph
        elif idx == _FOOTER_PLACEHOLDER_IDX and ph.has_text_frame:
            footer_ph_slots.append(ph)
        elif ph.has_text_frame:
            body_phs.append(ph)

    # If there are no non-footer body placeholders, fall back to the footer slot
    # (e.g. Statement layout where idx=10 is the only body area)
    if not body_phs:
        body_phs = footer_ph_slots

    if title_ph is None:
        return False

    # Fill title
    title_ph.text_frame.clear()
    title_ph.text_frame.paragraphs[0].text = title

    # Sort remaining body placeholders top→bottom, then left→right
    body_phs.sort(key=lambda p: (p.top or 0, p.left or 0))

    def _fill_ph(ph, lines: list) -> None:
        tf = ph.text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line

    if body and body_phs:
        _fill_ph(body_phs[0], body)
    if body2 and len(body_phs) >= 2:
        _fill_ph(body_phs[1], body2)

    return True


def _blank(layouts: dict):
    """Best blank-style layout for slides we fully control."""
    return layouts.get("blank", layouts.get("title_only", layouts["content"]))


# ─── Slide Builders ───────────────────────────────────────────────────────────

def _build_title_slide(slide, structure: PresentationStructure,
                       g: SlideGeometry):
    """
    Title slide — left-aligned layout with thin accent rule.
    Template placeholders removed; all elements drawn at fixed positions.
    """
    _remove_all_placeholders(slide)
    _add_top_bar(slide, g)

    # Main title
    size = min(max(_title_size(structure.title, 40), 30), 48)
    ty   = 0.28
    txb  = _textbox(slide, g.x(0.08), g.y(ty), g.w(0.84), g.h(0.22))
    _vcenter(txb)
    _enable_auto_shrink(txb)
    p    = txb.text_frame.paragraphs[0]
    _set_para(p, structure.title, size, bold=True,
              color=COLORS["text_dark"], align=PP_ALIGN.LEFT)

    # Thin accent rule below title
    rule_y = ty + 0.23
    _rect(slide, g.x(0.08), g.y(rule_y), g.w(0.10), g.h(0.003), COLORS["navy"])

    # Subtitle
    sub_y = rule_y + 0.015
    sub   = _textbox(slide, g.x(0.08), g.y(sub_y), g.w(0.65), g.h(0.12))
    _enable_auto_shrink(sub)
    p2    = sub.text_frame.paragraphs[0]
    _set_para(p2, structure.subtitle, 22, color=COLORS["text_body"])

    # Author / date at bottom
    if structure.author:
        auth = _textbox(slide, g.x(ML), g.y(0.88), g.w(CW), g.h(0.06))
        p3   = auth.text_frame.paragraphs[0]
        _set_para(p3, structure.author, 14, color=COLORS["text_muted"])


def _build_section_slide(
    slide, content: SlideContent, g: SlideGeometry, use_native: bool = False
):
    """Section-divider slide — large left accent bar + large centered title.

    When *use_native* is True the slide was added with a native template layout
    (e.g. Nagarro's blue/green section layout).  In that case we fill the
    layout's own placeholders and skip all manual drawing so the template's
    original design shows through.
    """
    if use_native:
        desc = [content.bullets[0]] if content.bullets else None
        _fill_native_placeholders(slide, content.title, desc)
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes
        return

    _remove_all_placeholders(slide)
    _add_top_bar(slide, g)

    # Vertical navy accent bar on the left
    _rect(slide, g.x(ML), g.y(0.28), g.w(0.007), g.h(0.44), COLORS["navy"])

    # Large section title
    text_x = g.x(ML) + g.w(0.007) + g.w(0.025)
    text_w = g.w(CW) - g.w(0.007) - g.w(0.025)
    size   = _title_size(content.title, 40)
    txb    = _textbox(slide, text_x, g.y(0.28), text_w, g.h(0.42))
    _vcenter(txb)
    _enable_auto_shrink(txb)
    p      = txb.text_frame.paragraphs[0]
    _set_para(p, content.title, size, bold=True, color=COLORS["text_dark"])

    # Optional description line
    if content.bullets:
        desc = _textbox(slide, text_x, g.y(0.72), text_w, g.h(0.14))
        p2   = desc.text_frame.paragraphs[0]
        _set_para(p2, content.bullets[0], 20, color=COLORS["text_muted"])

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _build_content_slide(slide, content: SlideContent, g: SlideGeometry):
    """
    Standard content slide — layout chosen automatically for visual variety:

    • 2–3 short non-emoji bullets  → side-by-side cards (pillars / options)
    • 2–5 emoji-prefixed bullets   → icon-list (pill + text rows)
    • 5+ bullets or long text      → traditional em-dash list
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    if content.bullets:
        if _use_card_layout(content.bullets):
            _build_bullet_cards(slide, g, content.bullets)
        elif _has_emoji_prefix(content.bullets) and len(content.bullets) <= 5:
            _build_icon_list(slide, g, content.bullets)
        else:
            size = _bullet_size(content.bullets)
            txb  = _textbox(slide, g.x(ML), g.y(g.body_y), g.w(CW), g.h(g.body_h))
            _vcenter(txb)
            _enable_auto_shrink(txb)
            _add_bullets(txb.text_frame, content.bullets,
                         size=size, bold_first_word=True)

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _build_chart_slide(slide, content: SlideContent, g: SlideGeometry):
    """
    Chart slide.
    With bullets  → chart left 60 %, insight card right 36 %.
    Without bullets → chart fills full content zone.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    if not content.charts:
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes
        return

    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)

    if content.bullets:
        chart_w = g.w(0.60)
        gap     = g.w(0.030)
        card_x  = bl + chart_w + gap
        card_w  = bw - chart_w - gap

        # Chart (left 60 %)
        _render_chart(slide, content.charts[0], bl, bt, chart_w, bh)

        # Insight card background + left accent
        _round_rect(slide, card_x, bt, card_w, bh, COLORS["bg_card"])
        _rect(slide, card_x, bt, g.w(0.005), bh, COLORS["blue"])

        inner_x = card_x + g.w(0.005) + g.w(0.015)
        inner_w = card_w - g.w(0.005) - g.w(0.030)
        pad_top = g.h(0.025)

        # "KEY INSIGHTS" label
        lbl = _textbox(slide, inner_x, bt + pad_top, inner_w, g.h(0.040))
        r   = lbl.text_frame.paragraphs[0].add_run()
        r.text           = "KEY INSIGHTS"
        r.font.name      = FONT
        r.font.size      = Pt(10)
        r.font.bold      = True
        r.font.color.rgb = COLORS["blue"]

        # Thin rule under label
        _rect(slide, inner_x, bt + pad_top + g.h(0.045),
              inner_w, g.h(0.002), COLORS["divider"])

        # Bullets in card
        bul_top = bt + pad_top + g.h(0.055)
        bul_h   = bh - pad_top - g.h(0.060)
        size    = min(_bullet_size(content.bullets), 14)
        btxb    = _textbox(slide, inner_x, bul_top, inner_w, bul_h)
        _vcenter(btxb)
        _enable_auto_shrink(btxb)
        _add_bullets(btxb.text_frame, content.bullets,
                     size=size, bold_first_word=True, line_spacing=1.2)
    else:
        _render_chart(slide, content.charts[0], bl, bt, bw, bh)

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _build_multi_chart_slide(slide, content: SlideContent, g: SlideGeometry):
    """Dashboard slide: 1–6 charts arranged in a gap-aware grid, rendered in parallel."""
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    charts = content.charts
    n      = len(charts)
    if n == 0:
        return

    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)
    gx = g.w(0.020)   # horizontal gap
    gy = g.h(0.025)   # vertical gap

    # ── Build slot list: (left, top, width, height) ───────────────────────────
    if n == 1:
        slots = [(bl, bt, bw, bh)]

    elif n == 2:
        cw = (bw - gx) // 2
        slots = [
            (bl,           bt, cw, bh),
            (bl + cw + gx, bt, cw, bh),
        ]

    elif n == 3:
        cw = (bw - 2 * gx) // 3
        slots = [(bl + i * (cw + gx), bt, cw, bh) for i in range(3)]

    elif n == 4:
        cw = (bw - gx) // 2
        ch = (bh - gy) // 2
        slots = [
            (bl,           bt,           cw, ch),
            (bl + cw + gx, bt,           cw, ch),
            (bl,           bt + ch + gy, cw, ch),
            (bl + cw + gx, bt + ch + gy, cw, ch),
        ]

    else:   # 5–6: 3 × 2 grid
        cw = (bw - 2 * gx) // 3
        ch = (bh - gy) // 2
        slots = [
            (bl + (i % 3) * (cw + gx), bt + (i // 3) * (ch + gy), cw, ch)
            for i in range(min(n, 6))
        ]

    charts_to_render = charts[:len(slots)]

    # ── Render all charts concurrently, then place sequentially ───────────────
    rendered = _render_charts_parallel(charts_to_render, slots)

    for img_bytes, chart_spec, (left, top, w, h) in zip(rendered, charts_to_render, slots):
        if img_bytes is not None:
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, w, h)
        else:
            txb = _textbox(slide, left, top, w, h)
            txb.text_frame.text = f"[Chart error: {chart_spec.chart_function}]"

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _build_key_number_slide(slide, content: SlideContent, g: SlideGeometry):
    """
    Hero-number slide.
    With chart   → chart top 55 %, bullets below.
    Without chart → KPI card centered in top 50 %, bullets below.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)

    if content.charts:
        chart_h    = int(bh * 0.55)
        _render_chart(slide, content.charts[0], bl, bt, bw, chart_h)
        bullet_top = bt + chart_h + g.h(0.025)
        bullet_h   = bh - chart_h - g.h(0.025)

    else:
        card_h = int(bh * 0.50)
        card_x = bl + g.w(0.20)
        card_w = g.w(0.50)

        _round_rect(slide, card_x, bt, card_w, card_h, COLORS["bg_accent"])
        _rect(slide, card_x, bt, g.w(0.005), card_h, COLORS["navy"])

        inner_x = card_x + g.w(0.005) + g.w(0.015)
        inner_w = card_w - g.w(0.005) - g.w(0.030)

        # Big number (top 52 % of card)
        num_h   = int(card_h * 0.52)
        num_txb = _textbox(slide, inner_x, bt + g.h(0.025), inner_w, num_h)
        p       = num_txb.text_frame.paragraphs[0]
        num_size = 60 if len(content.key_number) <= 6 else 48
        _set_para(p, content.key_number, num_size,
                  bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)

        # Label (below number, no overlap)
        lbl_top = bt + g.h(0.025) + num_h + g.h(0.008)
        lbl_h   = card_h - g.h(0.025) - num_h - g.h(0.008) - g.h(0.020)
        lbl_txb = _textbox(slide, inner_x, lbl_top, inner_w, lbl_h)
        p2      = lbl_txb.text_frame.paragraphs[0]
        _set_para(p2, content.key_label, 16,
                  color=COLORS["text_body"], align=PP_ALIGN.CENTER)

        bullet_top = bt + card_h + g.h(0.030)
        bullet_h   = bh - card_h - g.h(0.030)

    if content.bullets and bullet_h > g.h(0.06):
        size = _bullet_size(content.bullets)
        btxb = _textbox(slide, bl, bullet_top, bw, bullet_h)
        _enable_auto_shrink(btxb)
        _add_bullets(btxb.text_frame, content.bullets, size=size)
        for p in btxb.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _build_two_column_slide(
    slide, content: SlideContent, g: SlideGeometry, use_native: bool = False
):
    """Two-column comparison: equal-width cards separated by a gap.

    When *use_native* is True the slide uses the template's native two-column
    layout — content fills its placeholders and the template theme is preserved.
    """
    if use_native:
        left  = ([content.left_heading]  if content.left_heading  else []) + (content.bullets       or [])
        right = ([content.right_heading] if content.right_heading else []) + (content.right_bullets or [])
        _fill_native_placeholders(slide, content.title, left or None, right or None)
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes
        return

    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    bl  = g.x(ML)
    bt  = g.y(g.body_y)
    bw  = g.w(CW)
    bh  = g.h(g.body_h)
    gap = g.w(0.030)
    cw  = (bw - gap) // 2

    _build_column_card(slide, g,
                       bl,            bt, cw, bh,
                       content.left_heading, content.bullets,
                       COLORS["blue"])
    _build_column_card(slide, g,
                       bl + cw + gap, bt, cw, bh,
                       content.right_heading, content.right_bullets,
                       COLORS["teal"])

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _build_column_card(slide, g: SlideGeometry,
                       left, top, col_w, col_h,
                       heading: str, bullets: list, accent: RGBColor):
    """Single column card: background + accent bar + heading + divider + bullets."""
    _round_rect(slide, left, top, col_w, col_h, COLORS["bg_card"])

    bar_w = g.w(0.005)
    _rect(slide, left, top, bar_w, col_h, accent)

    pad_x   = g.w(0.018)
    pad_y   = g.h(0.025)
    inner_x = left + bar_w + pad_x
    inner_w = col_w - bar_w - 2 * pad_x

    # Heading
    heading_h = g.h(0.075)
    if heading:
        h_txb = _textbox(slide, inner_x, top + pad_y, inner_w, heading_h)
        _vcenter(h_txb)
        _enable_auto_shrink(h_txb)
        p     = h_txb.text_frame.paragraphs[0]
        _set_para(p, heading, 18, bold=True, color=accent)

    # Thin rule below heading
    rule_top = top + pad_y + heading_h + g.h(0.005)
    _rect(slide, inner_x, rule_top, inner_w, g.h(0.002), COLORS["divider"])

    # Bullets — guaranteed to not exceed card boundary
    if bullets:
        bul_top = rule_top + g.h(0.015)
        bul_h   = (top + col_h) - bul_top - pad_y
        if bul_h > g.h(0.04):
            size = min(_bullet_size(bullets), 16)
            btxb = _textbox(slide, inner_x, bul_top, inner_w, bul_h)
            _vcenter(btxb)
            _enable_auto_shrink(btxb)
            _add_bullets(btxb.text_frame, bullets,
                         size=size, bold_first_word=True,
                         color=COLORS["text_body"],
                         line_spacing=1.15)


# ─── Agenda Slide ─────────────────────────────────────────────────────────────

def _build_agenda_slide(slide, content: SlideContent, g: SlideGeometry) -> None:
    """
    Agenda / table-of-contents slide.

    Design adapts to item count:
      2–4 bullets → horizontal numbered cards (portrait, side by side).
      5–6 bullets → two-column numbered list (circle icon + title + desc).

    Bullet format:
      "Section Title"               → title only
      "Section Title: Description"  → title with muted description below

    Args:
        slide:   Target slide object.
        content: SlideContent; bullets = agenda items.
        g:       SlideGeometry for this presentation.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title or "Agenda")

    bullets = content.bullets
    if not bullets:
        return

    n  = len(bullets)
    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)

    def _parse(b: str) -> tuple[str, str]:
        if ":" in b:
            t, d = b.split(":", 1)
            return t.strip(), d.strip()
        return b.strip(), ""

    # ── Layout A: horizontal cards (2–4 items) ─────────────────────────────────
    if n <= 4:
        gx = g.w(0.022)
        cw = (bw - (n - 1) * gx) // n
        ch = bh

        num_size  = 52 if n <= 3 else 44
        ttl_size  = 18 if n <= 3 else 16
        desc_size = 12 if n <= 3 else 11

        for i, bullet in enumerate(bullets):
            title, desc = _parse(bullet)
            cx = bl + i * (cw + gx)

            # Card background
            _round_rect(slide, cx, bt, cw, ch, COLORS["bg_card"])

            # Top accent stripe (full card width)
            _rect(slide, cx, bt, cw, g.h(0.010), COLORS["navy"])

            # Large number
            num_h   = g.h(0.165)
            num_txb = _textbox(slide, cx + g.w(0.018), bt + g.h(0.040),
                               cw - g.w(0.036), num_h)
            num_p            = num_txb.text_frame.paragraphs[0]
            num_p.font.name  = FONT
            num_p.font.size  = Pt(num_size)
            num_p.font.bold  = True
            num_p.font.color.rgb = COLORS["navy"]
            num_p.text       = f"{i + 1:02d}"

            # Thin rule below number
            rule_y = bt + g.h(0.040) + num_h + g.h(0.008)
            _rect(slide, cx + g.w(0.018), rule_y,
                  cw - g.w(0.036), g.h(0.002), COLORS["divider"])

            # Section title
            title_y = rule_y + g.h(0.020)
            title_h = g.h(0.14)
            t_txb   = _textbox(slide, cx + g.w(0.018), title_y,
                               cw - g.w(0.036), title_h)
            t_txb.text_frame.word_wrap = True
            _enable_auto_shrink(t_txb)
            t_p            = t_txb.text_frame.paragraphs[0]
            t_p.font.name  = FONT
            t_p.font.size  = Pt(ttl_size)
            t_p.font.bold  = True
            t_p.font.color.rgb = COLORS["text_dark"]
            t_p.text       = title

            # Description (optional)
            if desc:
                desc_y = title_y + title_h + g.h(0.008)
                desc_h = (bt + ch) - desc_y - g.h(0.020)
                if desc_h >= g.h(0.040):
                    d_txb = _textbox(slide, cx + g.w(0.018), desc_y,
                                     cw - g.w(0.036), desc_h)
                    d_txb.text_frame.word_wrap = True
                    _enable_auto_shrink(d_txb)
                    d_p            = d_txb.text_frame.paragraphs[0]
                    d_p.font.name  = FONT
                    d_p.font.size  = Pt(desc_size)
                    d_p.font.color.rgb = COLORS["text_muted"]
                    d_p.text       = desc

    # ── Layout B: two-column numbered list (5–6 items) ─────────────────────────
    else:
        cols    = 2
        col_w   = (bw - g.w(0.040)) // cols
        rows    = (n + 1) // 2
        row_h   = bh // rows
        circ_d  = g.h(0.065)   # number circle diameter

        for i, bullet in enumerate(bullets):
            title, desc = _parse(bullet)
            col = i % cols
            row = i // cols
            rx  = bl + col * (col_w + g.w(0.040))
            ry  = bt + row * row_h

            # Horizontal rule above each row (except first)
            if row > 0 and col == 0:
                _rect(slide, bl, ry - g.h(0.010), bw, g.h(0.001), COLORS["divider"])

            # Number circle (vertically centred in row)
            circ_y = ry + (row_h - circ_d) // 2
            _round_rect(slide, rx, circ_y, circ_d, circ_d, COLORS["navy"])
            c_txb = _textbox(slide, rx, circ_y, circ_d, circ_d)
            _vcenter(c_txb)
            cp            = c_txb.text_frame.paragraphs[0]
            cp.font.name  = FONT
            cp.font.size  = Pt(14)
            cp.font.bold  = True
            cp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cp.alignment  = PP_ALIGN.CENTER
            cp.text       = f"{i + 1:02d}"

            # Text block: title + description
            tx_x = rx + circ_d + g.w(0.015)
            tx_w = col_w - circ_d - g.w(0.015)

            half = row_h // 2
            t_txb = _textbox(slide, tx_x, ry + g.h(0.012), tx_w, half)
            t_txb.text_frame.word_wrap = True
            _enable_auto_shrink(t_txb)
            tp            = t_txb.text_frame.paragraphs[0]
            tp.font.name  = FONT
            tp.font.size  = Pt(15)
            tp.font.bold  = True
            tp.font.color.rgb = COLORS["text_dark"]
            tp.text       = title

            if desc:
                d_txb = _textbox(slide, tx_x, ry + half, tx_w,
                                 row_h - half - g.h(0.012))
                d_txb.text_frame.word_wrap = True
                _enable_auto_shrink(d_txb)
                dp            = d_txb.text_frame.paragraphs[0]
                dp.font.name  = FONT
                dp.font.size  = Pt(12)
                dp.font.color.rgb = COLORS["text_muted"]
                dp.text       = desc

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


# ─── Metrics Grid Slide ───────────────────────────────────────────────────────

def _build_metrics_grid_slide(slide, content: SlideContent, g: SlideGeometry) -> None:
    """
    Manus-style KPI metrics grid with left-accent cards.

    Items field format (required):
      [{"value": "–28%", "label": "Operative Kosten", "trend": "down"},
       {"value": "+35%", "label": "Produktivität",    "trend": "up"}, ...]

    Trend values: "up" → ▲  |  "down" → ▼  |  omit for no arrow.

    If key_number is also set, a dark hero panel is rendered on the right:
      • key_label  → panel heading (e.g. "Netto-ROI Jahr 1")
      • key_number → large hero number (e.g. "567 %") in accent colour
      • bullets    → optional calculation rows ("Einsparung: 480.000 €")

    Grid sizing:
      ≤ 4 items → 2 × 2   |   ≤ 6 items → 2 × 3

    Args:
        slide:   Target slide object.
        content: SlideContent with items list and optional key_number / bullets.
        g:       SlideGeometry for this presentation.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    items = content.items
    if not items:
        return

    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)

    has_hero = bool(content.key_number)
    grid_w   = g.w(0.565) if has_hero else bw
    hero_gap = g.w(0.028)
    hero_x   = bl + grid_w + hero_gap
    hero_w   = bw - grid_w - hero_gap

    # ── Metric cards grid ──────────────────────────────────────────────────────
    cols = 2
    rows = max(1, (len(items[:6]) + 1) // 2)
    gx   = g.w(0.020)
    gy   = g.h(0.022)
    cw   = (grid_w - (cols - 1) * gx) // cols
    ch   = (bh - (rows - 1) * gy) // rows

    ACCENT_BAR_W = g.w(0.007)   # Manus-style left accent border width

    for i, item in enumerate(items[: cols * rows]):
        col = i % cols
        row = i // cols
        cx  = bl + col * (cw + gx)
        cy  = bt + row * (ch + gy)

        value = str(item.get("value", ""))
        label = str(item.get("label", ""))
        trend = str(item.get("trend", ""))

        # Card background + left accent bar (Manus pattern)
        _round_rect(slide, cx, cy, cw, ch, COLORS["bg_card"])
        _rect(slide, cx, cy, ACCENT_BAR_W, ch, COLORS["navy"])

        inner_x = cx + ACCENT_BAR_W + g.w(0.012)
        inner_w = cw - ACCENT_BAR_W - g.w(0.018)

        # Big metric value
        v_size   = 36 if len(value) <= 4 else 30 if len(value) <= 7 else 24
        val_top  = cy + int(ch * 0.14)
        val_h    = int(ch * 0.44)
        val_txb  = _textbox(slide, inner_x, val_top, inner_w, val_h)
        _enable_auto_shrink(val_txb)
        val_p    = val_txb.text_frame.paragraphs[0]
        val_p.font.name      = FONT
        val_p.font.size      = Pt(v_size)
        val_p.font.bold      = True
        val_p.font.color.rgb = COLORS["navy"]
        val_p.text           = value

        # Trend arrow as a second run
        if trend == "up":
            trend_run       = val_p.add_run()
            trend_run.text  = " ▲"
            trend_run.font.name      = FONT
            trend_run.font.size      = Pt(max(14, v_size - 12))
            trend_run.font.color.rgb = COLORS["navy"]
        elif trend == "down":
            trend_run       = val_p.add_run()
            trend_run.text  = " ▼"
            trend_run.font.name      = FONT
            trend_run.font.size      = Pt(max(14, v_size - 12))
            trend_run.font.color.rgb = COLORS["navy"]

        # Metric label below value
        lbl_top = cy + int(ch * 0.58)
        lbl_h   = int(ch * 0.30)
        lbl_txb = _textbox(slide, inner_x, lbl_top, inner_w, lbl_h)
        _enable_auto_shrink(lbl_txb)
        lbl_p   = lbl_txb.text_frame.paragraphs[0]
        lbl_p.font.name      = FONT
        lbl_p.font.size      = Pt(12)
        lbl_p.font.color.rgb = COLORS["text_muted"]
        lbl_p.text           = label

    if not has_hero:
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes
        return

    # ── Dark hero panel ────────────────────────────────────────────────────────
    # Use a medium-dark slate (#2D3748) — noticeably darker than white-bg slides
    # yet clearly lighter than Nagarro's very-dark bg (#06041F).  Works on all
    # templates without needing to know the exact background hue.
    DARK_FILL = RGBColor(0x2D, 0x37, 0x48)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED_W   = RGBColor(0xA0, 0xA0, 0xB8)
    DIVIDER_D = RGBColor(0x44, 0x44, 0x60)

    _round_rect(slide, hero_x, bt, hero_w, bh, DARK_FILL)

    pad     = g.w(0.025)
    inner_x = hero_x + pad
    inner_w = hero_w - 2 * pad

    # ── Panel heading ──────────────────────────────────────────────────────────
    hero_label = content.key_label or "ROI"
    hl_txb = _textbox(slide, inner_x, bt + g.h(0.030), inner_w, g.h(0.060))
    hl_p   = hl_txb.text_frame.paragraphs[0]
    hl_p.font.name      = FONT
    hl_p.font.size      = Pt(12)
    hl_p.font.bold      = True
    hl_p.font.color.rgb = MUTED_W
    hl_p.text           = hero_label.upper()

    # ── Calculation rows from bullets ──────────────────────────────────────────
    y_calc    = bt + g.h(0.105)
    row_h_c   = g.h(0.052)
    _rect(slide, inner_x, y_calc - g.h(0.008), inner_w, g.h(0.002), DIVIDER_D)

    for bullet in content.bullets:
        calc_txb = _textbox(slide, inner_x, y_calc, inner_w, row_h_c)
        cp       = calc_txb.text_frame.paragraphs[0]
        if ":" in bullet:
            lbl_part, val_part = bullet.split(":", 1)
            r1 = cp.add_run()
            r1.text           = lbl_part.strip()
            r1.font.name      = FONT
            r1.font.size      = Pt(12)
            r1.font.color.rgb = MUTED_W
            r2 = cp.add_run()
            r2.text           = ":  " + val_part.strip()
            r2.font.name      = FONT
            r2.font.size      = Pt(12)
            r2.font.bold      = True
            r2.font.color.rgb = WHITE
        else:
            cp.font.name      = FONT
            cp.font.size      = Pt(12)
            cp.font.color.rgb = WHITE
            cp.text           = bullet
        y_calc += row_h_c

    # Divider before big number
    _rect(slide, inner_x, y_calc + g.h(0.010), inner_w, g.h(0.002), DIVIDER_D)

    # ── Big hero number ────────────────────────────────────────────────────────
    num      = content.key_number
    n_size   = 58 if len(num) <= 4 else 46 if len(num) <= 6 else 36
    num_top  = y_calc + g.h(0.040)
    num_h    = (bt + bh) - num_top - g.h(0.080)
    num_txb  = _textbox(slide, inner_x, num_top, inner_w, num_h)
    _vcenter(num_txb)
    _enable_auto_shrink(num_txb)
    num_p    = num_txb.text_frame.paragraphs[0]
    num_p.font.name      = FONT
    num_p.font.size      = Pt(n_size)
    num_p.font.bold      = True
    num_p.font.color.rgb = COLORS["navy"]   # template accent (teal on Nagarro)
    num_p.alignment      = PP_ALIGN.CENTER
    num_p.text           = num

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


# ─── Pricing Slide ────────────────────────────────────────────────────────────

def _build_pricing_slide(slide, content: SlideContent, g: SlideGeometry) -> None:
    """
    3-tier pricing cards — Manus pricing-slide pattern.

    Items field format (2–3 entries):
      [
        {"tier": "Growth",       "price": "ab 1.500 €", "period": "pro Monat",
         "target": "KMU",        "features": ["Bis zu 50 Nutzer", "3 Module"]},
        {"tier": "Professional", "price": "ab 4.500 €", "period": "pro Monat",
         "target": "Mittelstand","features": ["250 Nutzer", "5 Module", "API"],
         "recommended": true},
        {"tier": "Enterprise",   "price": "Individuell","period": "auf Anfrage",
         "target": "Konzerne",   "features": ["Unbegrenzte Nutzer", "Alle Module"],
         "dark": true}
      ]

    Design rules:
      • First tier  → neutral top stripe, white card.
      • recommended → accent colour stripe, slightly taller card, "✦ Empfohlen" badge.
      • dark        → dark-navy card, teal stripe, white text.

    Args:
        slide:   Target slide object.
        content: SlideContent with items list of tier dicts.
        g:       SlideGeometry for this presentation.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    items   = content.items
    n_tiers = min(len(items), 3)
    if n_tiers < 2:
        return

    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)
    gx = g.w(0.025)
    cw = (bw - (n_tiers - 1) * gx) // n_tiers

    DARK_FILL   = RGBColor(0x1A, 0x20, 0x2C)   # deep slate — works on all templates
    DARK_STRIPE = _DEFAULT_COLORS["teal"]
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED_W     = RGBColor(0xA0, 0xB0, 0xC0)
    CHECK_GREEN = RGBColor(0x48, 0xBB, 0x78)
    NEUTRAL_STR = COLORS["divider"]

    for i, tier in enumerate(items[:3]):
        cx = bl + i * (cw + gx)

        is_recommended = bool(tier.get("recommended", False))
        is_dark        = bool(tier.get("dark", False))

        tier_name   = str(tier.get("tier",    f"Tier {i + 1}"))
        tier_price  = str(tier.get("price",   ""))
        tier_period = str(tier.get("period",  ""))
        tier_target = str(tier.get("target",  ""))
        features    = list(tier.get("features", []))

        # Recommended card is taller (elevated top + bottom)
        elev   = g.h(0.018) if is_recommended else 0
        card_t = bt - elev
        card_h = bh + 2 * elev

        # ── Card background ────────────────────────────────────────────────────
        if is_dark:
            _round_rect(slide, cx, card_t, cw, card_h, DARK_FILL)
            text_c  = WHITE
            muted_c = MUTED_W
        else:
            _round_rect(slide, cx, card_t, cw, card_h, WHITE)
            text_c  = COLORS["text_dark"]
            muted_c = COLORS["text_muted"]

        # ── Top colour stripe ──────────────────────────────────────────────────
        if is_recommended:
            stripe_c = COLORS["navy"]
        elif is_dark:
            stripe_c = DARK_STRIPE
        else:
            stripe_c = NEUTRAL_STR
        _rect(slide, cx, card_t, cw, g.h(0.010), stripe_c)

        # ── "Empfohlen" badge (recommended tier only) ──────────────────────────
        if is_recommended:
            bdg_w = g.w(0.100)
            bdg_h = g.h(0.038)
            bdg_x = cx + (cw - bdg_w) // 2
            bdg_y = card_t - bdg_h // 2
            _round_rect(slide, bdg_x, bdg_y, bdg_w, bdg_h, COLORS["navy"])
            bdg_txb = _textbox(slide, bdg_x, bdg_y, bdg_w, bdg_h)
            _vcenter(bdg_txb)
            bdg_p            = bdg_txb.text_frame.paragraphs[0]
            bdg_p.font.name  = FONT
            bdg_p.font.size  = Pt(10)
            bdg_p.font.bold  = True
            bdg_p.font.color.rgb = WHITE
            bdg_p.alignment  = PP_ALIGN.CENTER
            bdg_p.text       = "✦  Empfohlen"

        # ── Card content ───────────────────────────────────────────────────────
        pad   = g.w(0.018)
        inner_x = cx + pad
        inner_w = cw - 2 * pad
        y       = card_t + g.h(0.038)   # cursor

        def _tier_txt(text: str, size: int, bold: bool = False, color: RGBColor = None):
            nonlocal y
            h    = g.h(0.055)
            txb  = _textbox(slide, inner_x, y, inner_w, h)
            para = txb.text_frame.paragraphs[0]
            para.font.name      = FONT
            para.font.size      = Pt(size)
            para.font.bold      = bold
            para.font.color.rgb = color or text_c
            para.text           = text
            y += h

        # Tier name (coloured)
        name_color = (stripe_c if not is_dark else RGBColor(0x4F, 0xD1, 0xC5))
        _tier_txt(tier_name.upper(), 15, bold=True, color=name_color)
        y -= g.h(0.010)  # tighten

        if tier_target:
            _tier_txt(tier_target, 11, color=muted_c)
            y -= g.h(0.008)

        # Price (large)
        if tier_price:
            price_h = g.h(0.070)
            ptxb    = _textbox(slide, inner_x, y, inner_w, price_h)
            pp_p    = ptxb.text_frame.paragraphs[0]
            pp_p.font.name      = FONT
            pp_p.font.size      = Pt(24)
            pp_p.font.bold      = True
            pp_p.font.color.rgb = text_c
            pp_p.text           = tier_price
            y += price_h

        if tier_period:
            _tier_txt(tier_period, 11, color=muted_c)
            y -= g.h(0.005)

        # Thin divider
        divider_c = RGBColor(0x44, 0x44, 0x60) if is_dark else COLORS["divider"]
        _rect(slide, inner_x, y + g.h(0.004), inner_w, g.h(0.002), divider_c)
        y += g.h(0.020)

        # Features with green checkmarks
        feat_row_h = min(g.h(0.042), max(g.h(0.030),
                         (card_t + card_h - y - g.h(0.020)) // max(1, len(features))))
        for feat in features:
            if y + feat_row_h > card_t + card_h - g.h(0.010):
                break   # prevent overflow
            ftxb = _textbox(slide, inner_x, y, inner_w, feat_row_h)
            _enable_auto_shrink(ftxb)
            fp   = ftxb.text_frame.paragraphs[0]
            r1   = fp.add_run()
            r1.text           = "✓  "
            r1.font.name      = FONT
            r1.font.size      = Pt(11)
            r1.font.color.rgb = CHECK_GREEN
            r2   = fp.add_run()
            r2.text           = feat
            r2.font.name      = FONT
            r2.font.size      = Pt(11)
            r2.font.color.rgb = text_c
            y += feat_row_h

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


# ─── Icon Grid Slide ──────────────────────────────────────────────────────────

def _build_icon_grid_slide(slide, content: SlideContent, g: SlideGeometry) -> None:
    """
    Icon grid layout: 3–6 feature cards arranged in a responsive grid.

    Each card has a coloured icon circle (emoji), bold title, and short
    description.  Ideal for capability overviews, benefit lists, or feature
    summaries.  Bullet format expected:
      "🚀 Feature Name: Short description of this feature"

    Grid sizing:
      3 items → 3 × 1 (single wide row)
      4 items → 2 × 2
      5–6 items → 3 × 2

    Args:
        slide:   Target slide object.
        content: SlideContent with bullets in emoji+title:desc format.
        g:       SlideGeometry for this presentation.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    bullets = content.bullets
    if not bullets:
        return

    n  = len(bullets)
    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)
    gx = g.w(0.020)
    gy = g.h(0.025)

    # ── Grid dimensions ────────────────────────────────────────────────────────
    if n <= 3:
        cols, rows = n, 1
    elif n == 4:
        cols, rows = 2, 2
    else:
        cols, rows = 3, 2

    cw = (bw - (cols - 1) * gx) // cols
    ch = (bh - (rows - 1) * gy) // rows

    icon_d = min(g.w(0.055), g.h(0.075))   # icon circle diameter

    for i, bullet in enumerate(bullets[: cols * rows]):
        col = i % cols
        row = i // cols
        cx  = bl + col * (cw + gx)
        cy  = bt + row * (ch + gy)

        # Card background + top accent stripe
        _round_rect(slide, cx, cy, cw, ch, COLORS["bg_card"])
        stripe_h = g.h(0.007)
        _rect(slide, cx, cy, cw, stripe_h, COLORS["navy"])

        # Parse emoji / label / description from bullet text
        if bullet and ord(bullet[0]) > 0x1F00:
            parts = bullet.split(" ", 1)
            emoji = parts[0]
            body  = parts[1].strip() if len(parts) > 1 else ""
        else:
            emoji = "●"
            body  = bullet

        # ── Icon circle ────────────────────────────────────────────────────────
        icon_x = cx + (cw - icon_d) // 2
        icon_y = cy + stripe_h + g.h(0.018)
        _round_rect(slide, icon_x, icon_y, icon_d, icon_d, COLORS["navy"])
        em_txb = _textbox(slide, icon_x, icon_y, icon_d, icon_d)
        _vcenter(em_txb)
        em_p       = em_txb.text_frame.paragraphs[0]
        em_p.text  = emoji
        em_p.font.size  = Pt(20 if rows == 1 else 18)
        em_p.alignment  = PP_ALIGN.CENTER

        # ── Title + description ────────────────────────────────────────────────
        text_y   = icon_y + icon_d + g.h(0.012)
        text_h   = (cy + ch) - text_y - g.h(0.012)
        text_pad = g.w(0.012)
        tx       = _textbox(slide, cx + text_pad, text_y,
                            cw - 2 * text_pad, text_h)
        tx.text_frame.word_wrap = True
        _enable_auto_shrink(tx)

        if ":" in body:
            card_title, desc = body.split(":", 1)
            p = tx.text_frame.paragraphs[0]
            p.font.name      = FONT
            p.font.size      = Pt(14 if rows >= 2 else 16)
            p.font.bold      = True
            p.font.color.rgb = COLORS["text_dark"]
            p.alignment      = PP_ALIGN.CENTER
            p.text           = card_title.strip()
            p2 = tx.text_frame.add_paragraph()
            p2.space_before       = Pt(4)
            p2.font.name          = FONT
            p2.font.size          = Pt(11 if rows >= 2 else 12)
            p2.font.color.rgb     = COLORS["text_muted"]
            p2.alignment          = PP_ALIGN.CENTER
            p2.text               = desc.strip()
        else:
            p = tx.text_frame.paragraphs[0]
            p.font.name      = FONT
            p.font.size      = Pt(14)
            p.font.bold      = True
            p.font.color.rgb = COLORS["text_dark"]
            p.alignment      = PP_ALIGN.CENTER
            p.text           = body

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


# ─── Timeline Slide ───────────────────────────────────────────────────────────

def _build_timeline_slide(slide, content: SlideContent, g: SlideGeometry) -> None:
    """
    Horizontal timeline: connecting rule + numbered/icon circles + step labels.

    Best for 3–5 sequential process steps or milestones.  Bullet format:
      "1️⃣ Phase Name: Brief description of what happens here"
    Or plain:
      "Phase Name: Description"  (auto-numbered with step index)

    Layout:
        ───●────────●────────●────────●───     ← connecting line + circles
          Step 1   Step 2   Step 3   Step 4
         Desc...  Desc...  Desc...  Desc...

    Args:
        slide:   Target slide object.
        content: SlideContent with bullets as timeline steps.
        g:       SlideGeometry for this presentation.
    """
    _remove_all_placeholders(slide)
    _slide_frame(slide, g, content.title)

    bullets = content.bullets
    if not bullets:
        return

    n  = len(bullets)
    bl = g.x(ML)
    bt = g.y(g.body_y)
    bw = g.w(CW)
    bh = g.h(g.body_h)

    # ── Geometry: line sits in upper-middle of content zone ────────────────────
    dot_d   = g.h(0.078)              # circle diameter
    line_y  = bt + int(bh * 0.38)    # vertical midpoint of the dot
    step_w  = bw // n

    # Connecting line (spans all dots)
    line_x1 = bl + step_w // 2
    line_x2 = bl + bw - step_w // 2
    _rect(slide, line_x1, line_y + dot_d // 2 - g.h(0.003),
          line_x2 - line_x1, g.h(0.006), COLORS["divider"])

    for i, bullet in enumerate(bullets):
        cx_mid = bl + i * step_w + step_w // 2

        # Parse icon / step number and body
        if bullet and ord(bullet[0]) > 0x1F00:
            parts = bullet.split(" ", 1)
            icon  = parts[0]
            body  = parts[1].strip() if len(parts) > 1 else ""
        else:
            icon  = str(i + 1)
            body  = bullet

        # Alternate dot accent: odd steps use a slightly lighter tone
        dot_color = COLORS["navy"] if i % 2 == 0 else COLORS["blue"]

        # ── Step circle ────────────────────────────────────────────────────────
        dot_x = cx_mid - dot_d // 2
        dot_y = line_y
        _round_rect(slide, dot_x, dot_y, dot_d, dot_d, dot_color)

        dot_txb = _textbox(slide, dot_x, dot_y, dot_d, dot_d)
        _vcenter(dot_txb)
        dp            = dot_txb.text_frame.paragraphs[0]
        dp.text       = icon if (icon and ord(icon[0]) > 0x1F00) else str(i + 1)
        dp.font.name  = FONT
        dp.font.size  = Pt(15)
        dp.font.bold  = True
        dp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        dp.alignment  = PP_ALIGN.CENTER

        # ── Label below dot ────────────────────────────────────────────────────
        label_x = cx_mid - step_w // 2 + g.w(0.008)
        label_w = step_w - g.w(0.016)
        label_y = dot_y + dot_d + g.h(0.018)
        label_h = (bt + bh) - label_y - g.h(0.010)

        if ":" in body:
            step_title, step_desc = body.split(":", 1)
        else:
            step_title = body
            step_desc  = ""

        lbl_txb = _textbox(slide, label_x, label_y, label_w, label_h)
        lbl_txb.text_frame.word_wrap = True
        _enable_auto_shrink(lbl_txb)
        title_size = 14 if n <= 4 else 12
        desc_size  = 11 if n <= 4 else 10

        p = lbl_txb.text_frame.paragraphs[0]
        p.font.name      = FONT
        p.font.size      = Pt(title_size)
        p.font.bold      = True
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment      = PP_ALIGN.CENTER
        p.text           = step_title.strip()

        if step_desc.strip():
            p2 = lbl_txb.text_frame.add_paragraph()
            p2.space_before       = Pt(5)
            p2.font.name          = FONT
            p2.font.size          = Pt(desc_size)
            p2.font.color.rgb     = COLORS["text_muted"]
            p2.alignment          = PP_ALIGN.CENTER
            p2.text               = step_desc.strip()

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


# ─── Quote Slide ──────────────────────────────────────────────────────────────

def _build_quote_slide(
    slide, content: SlideContent, g: SlideGeometry, use_native: bool = False
) -> None:
    """
    Quote / highlight slide: large centered quotation with attribution.

    Maps schema fields:
      title      → optional context label shown in small caps above the quote
      key_number → the quote text (primary)
      key_label  → attribution (e.g. "CEO, Acme Corp")
      bullets    → fallback: bullets[0] = quote, bullets[1] = attribution

    Design: decorative large open-quote glyph, italic body, thin accent rule,
    bold attribution line.  Lots of white space — no bullets, no charts.

    When *use_native* is True the template's Statement/Quotes layout is used
    and its placeholders are filled — preserving the native visual design.

    Args:
        slide:   Target slide object.
        content: SlideContent with key_number (quote) and key_label (author).
        g:       SlideGeometry for this presentation.
    """
    quote_text  = content.key_number or (content.bullets[0] if content.bullets else "")
    attribution = content.key_label  or (content.bullets[1] if len(content.bullets) > 1 else "")

    if use_native:
        # In Statement layouts the TITLE placeholder is the large central text.
        # Put the quote there; attribution goes in the body/footer slot.
        quote_display = f"\u201C{quote_text}\u201D" if quote_text else content.title or ""
        attrib_lines  = [f"\u2014 {attribution}"] if attribution else []
        _fill_native_placeholders(slide, quote_display, attrib_lines or None)
        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes
        return

    _remove_all_placeholders(slide)
    _add_top_bar(slide, g)
    # No divider — maximise white space for the quote

    bl = g.x(ML)
    bw = g.w(CW)

    quote_text  = (content.key_number
                   or (content.bullets[0] if content.bullets else ""))
    attribution = (content.key_label
                   or (content.bullets[1] if len(content.bullets) > 1 else ""))

    # ── Decorative large open-quote glyph ─────────────────────────────────────
    qq_txb = _textbox(slide, g.x(ML), g.y(0.13), g.w(0.18), g.h(0.18))
    qq_p   = qq_txb.text_frame.paragraphs[0]
    qq_p.text           = "\u201C"
    qq_p.font.name      = FONT
    qq_p.font.size      = Pt(110)
    qq_p.font.bold      = True
    qq_p.font.color.rgb = COLORS["navy"]
    qq_p.alignment      = PP_ALIGN.LEFT

    # ── Context label (optional) ───────────────────────────────────────────────
    if content.title:
        ctx_txb = _textbox(slide, bl, g.y(0.20), bw, g.h(0.07))
        ctx_p   = ctx_txb.text_frame.paragraphs[0]
        ctx_p.font.name      = FONT
        ctx_p.font.size      = Pt(12)
        ctx_p.font.bold      = True
        ctx_p.font.color.rgb = COLORS["text_muted"]
        ctx_p.alignment      = PP_ALIGN.CENTER
        ctx_p.text           = content.title.upper()

    # ── Quote body text ────────────────────────────────────────────────────────
    q_size   = (28 if len(quote_text) <= 80
                else (22 if len(quote_text) <= 160 else 18))
    q_top    = 0.26
    q_txb    = _textbox(slide, bl + g.w(0.04), g.y(q_top),
                        bw - g.w(0.08), g.h(0.46))
    _vcenter(q_txb)
    q_txb.text_frame.word_wrap = True
    q_p = q_txb.text_frame.paragraphs[0]
    q_p.font.name      = FONT
    q_p.font.size      = Pt(q_size)
    q_p.font.italic    = True
    q_p.font.color.rgb = COLORS["text_dark"]
    q_p.alignment      = PP_ALIGN.CENTER
    q_p.text           = f"\u201C{quote_text}\u201D"
    _set_line_spacing(q_p, 1.45)

    # ── Thin accent rule before attribution ───────────────────────────────────
    _rect(slide, g.x(0.35), g.y(0.75), g.w(0.30), g.h(0.003), COLORS["navy"])

    # ── Attribution ───────────────────────────────────────────────────────────
    if attribution:
        auth_txb = _textbox(slide, bl, g.y(0.77), bw, g.h(0.09))
        auth_p   = auth_txb.text_frame.paragraphs[0]
        auth_p.font.name      = FONT
        auth_p.font.size      = Pt(16)
        auth_p.font.bold      = True
        auth_p.font.color.rgb = COLORS["text_muted"]
        auth_p.alignment      = PP_ALIGN.CENTER
        auth_p.text           = f"\u2014 {attribution}"

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


# ─── Main Entry Point ─────────────────────────────────────────────────────────

def generate_pptx(
    template_bytes: bytes,
    structure: PresentationStructure,
    template_colors: Optional[dict] = None,
    scraped_images: Optional[list] = None,
) -> bytes:
    """
    Assemble a fully-rendered PPTX from a template and AI-generated structure.

    Args:
        template_bytes:   Raw .pptx file bytes to use as the design base.
        structure:        Validated PresentationStructure from the AI service.
        template_colors:  Optional dict {bg, accent, text, muted} (hex strings)
                          used to derive the slide colour palette.
        scraped_images:   Optional list of Path objects to scraped web images.

    Returns:
        Raw bytes of the finished .pptx file.
    """
    global COLORS
    COLORS = _build_template_palette(template_colors) if template_colors else _DEFAULT_COLORS

    # Sync Plotly chart theme with the active template's accent color
    if template_colors and "accent" in template_colors:
        bg_hex  = template_colors.get("bg", "#FFFFFF")
        set_chart_theme(template_colors["accent"], is_dark=_is_dark_color(bg_hex))

    prs          = Presentation(io.BytesIO(template_bytes))
    logo_safe_y  = _detect_logo_safe_y(prs)
    layouts      = _discover_layouts(prs)
    g            = SlideGeometry(prs, logo_safe_y)
    blank        = _blank(layouts)

    # Remove all existing template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    total = len(structure.slides) + 1

    # ── Title slide ────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(layouts.get("title", blank))
    _build_title_slide(slide, structure, g)

    # ── Section layout rotation: alternate blue → green → white for variety ────
    _SECTION_CYCLE = ["section_blue", "section_green", "section_small_blue",
                      "section_white"]
    _section_idx = 0

    # ── Content slides ─────────────────────────────────────────────────────────
    for i, content in enumerate(structure.slides):
        lt  = content.layout_type
        num = i + 2

        # ── Choose layout ──────────────────────────────────────────────────────
        # Native layouts let the corporate template's own design show through.
        # Only simple slide types can safely use native layouts; chart-heavy
        # and custom-rendered slides always use blank.
        use_native    = False
        chosen_layout = blank

        if lt == "section_header":
            # Cycle through available section colour variants
            for candidate in _SECTION_CYCLE[_section_idx % len(_SECTION_CYCLE):] \
                             + _SECTION_CYCLE[:_section_idx % len(_SECTION_CYCLE)]:
                if candidate in layouts:
                    chosen_layout = layouts[candidate]
                    use_native    = True
                    break
            _section_idx += 1

        elif lt in ("two_column", "comparison"):
            native = layouts.get("two_column_native") or layouts.get("compare_native")
            if native:
                chosen_layout = native
                use_native    = True

        elif lt == "quote":
            native = (layouts.get("statement")
                      or layouts.get("quote_blue")
                      or layouts.get("quote_purple"))
            if native:
                chosen_layout = native
                use_native    = True

        slide = prs.slides.add_slide(chosen_layout)

        # ── Build slide ────────────────────────────────────────────────────────
        if lt == "section_header":
            _build_section_slide(slide, content, g, use_native=use_native)
        elif lt == "chart":
            _build_chart_slide(slide, content, g)
        elif lt == "multi_chart":
            _build_multi_chart_slide(slide, content, g)
        elif lt == "key_number":
            _build_key_number_slide(slide, content, g)
        elif lt in ("two_column", "comparison"):
            _build_two_column_slide(slide, content, g, use_native=use_native)
        elif lt == "icon_grid":
            _build_icon_grid_slide(slide, content, g)
        elif lt == "timeline":
            _build_timeline_slide(slide, content, g)
        elif lt == "quote":
            _build_quote_slide(slide, content, g, use_native=use_native)
        elif lt == "agenda":
            _build_agenda_slide(slide, content, g)
        elif lt == "metrics_grid":
            _build_metrics_grid_slide(slide, content, g)
        elif lt == "pricing":
            _build_pricing_slide(slide, content, g)
        else:
            _build_content_slide(slide, content, g)

        # Slide number only on manually-drawn slides (native layouts have their own)
        if not use_native:
            _add_slide_number(slide, g, num, total)

    # ── Insert scraped web images (appendix slide) ─────────────────────────────
    if scraped_images:
        from pathlib import Path as _Path
        valid_images = [p for p in scraped_images if _Path(p).exists()]
        if valid_images:
            # Add up to 4 images on a single appendix slide
            imgs_to_add = valid_images[:4]
            slide = prs.slides.add_slide(blank)
            _fill_slide_background(slide, g)
            _add_accent_bar(slide, g)

            # Title
            _add_textbox(
                slide, g.content_left, g.title_top, g.content_width, g.title_height,
                "Web Sources", COLORS["900"], Pt(22), bold=True,
            )

            # Grid layout for images
            cols = 2 if len(imgs_to_add) > 1 else 1
            rows = (len(imgs_to_add) + cols - 1) // cols
            img_w = g.content_width // cols - Inches(0.15)
            img_h = g.content_height // rows - Inches(0.15)
            for idx, img_path in enumerate(imgs_to_add):
                try:
                    col = idx % cols
                    row = idx // cols
                    left = g.content_left + col * (img_w + Inches(0.15))
                    top  = g.content_top + row * (img_h + Inches(0.15))
                    slide.shapes.add_picture(
                        str(img_path), left, top, img_w, img_h,
                    )
                except Exception as exc:
                    logger.error("Failed to insert scraped image %s: %s", img_path, exc)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
