"""Inspect generated PPTX to verify quality."""
from pptx import Presentation
from pptx.util import Inches, Pt
import sys

pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output_presentation.pptx"
prs = Presentation(pptx_path)

print(f"Slide dimensions: {prs.slide_width / 914400:.1f}\" x {prs.slide_height / 914400:.1f}\"")
print(f"Total slides: {len(prs.slides)}")
print("=" * 60)

for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f"\n--- Slide {i+1} (Layout: {layout_name}) ---")

    shapes = list(slide.shapes)
    print(f"  Shapes: {len(shapes)}")

    for shape in shapes:
        shape_type = shape.shape_type
        name = shape.name

        if shape.has_text_frame:
            text = shape.text_frame.text[:80]
            if text.strip():
                print(f"  [{name}] Text: {text}")

        if shape.has_chart:
            chart = shape.chart
            chart_type = str(chart.chart_type)
            series_count = len(chart.series)
            if chart.series:
                pts = len(chart.series[0].values)
                print(f"  [{name}] CHART: type={chart_type}, series={series_count}, data_points={pts}")

    # Check notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            print(f"  Speaker Notes: {notes[:60]}...")
